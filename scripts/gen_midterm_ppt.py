#!/usr/bin/env python3
"""
生成中期答辩 PPT（基于 中期报告_claude_输出.md 和 figures/ 中的图表）
运行方式：python scripts/gen_midterm_ppt.py
输出：中期答辩.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色方案 ────────────────────────────────────────────
BLUE   = RGBColor(0x1B, 0x3A, 0x5C)   # 深蓝主色
ACCENT = RGBColor(0x27, 0x7D, 0xA1)   # 蓝绿强调
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x33, 0x33, 0x33)
GRAY   = RGBColor(0x88, 0x88, 0x88)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FIG_DIR = os.path.join(os.path.dirname(__file__), '..', 'figures')
OUT_PATH = os.path.join(os.path.dirname(__file__), '..', '中期答辩.pptx')


def fig(name):
    p = os.path.join(FIG_DIR, name)
    return p if os.path.exists(p) else None


# ── 辅助函数 ────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, y=Inches(0.3), color=BLUE):
    """页面顶部深色标题条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), y - Inches(0.1),
                                 SLIDE_W, Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text_box(slide, left, top, width, height, lines,
                 font_size=16, bold=False, color=BLACK, spacing=1.2,
                 alignment=PP_ALIGN.LEFT, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(font_size * spacing * 0.4)
        p.alignment = alignment
        if bullet and line.strip():
            p.level = 0
    return txBox


def add_conclusion(slide, text, top=Inches(6.6)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), top,
                                  Inches(11.7), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_image_safe(slide, path, left, top, width=None, height=None):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    return False


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=15, color=BLACK, title=None, title_color=BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(font_size + 2)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(8)
    
    for i, item in enumerate(items):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return txBox


# ── 创建 PPT ────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout


# ═══════════════════════ P1 封面 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BLUE)

add_text_box(s, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
             ["基于零知识证明的分布式推理可验证性研究"],
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             ["中期答辩"],
             font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

info = ["学生：武垚乐", "指导教师：张宇超", "学院：计算机学院", "2026 年 3 月"]
add_text_box(s, Inches(3), Inches(4.8), Inches(7), Inches(2),
             info, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════ P2 背景与目标 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "课题背景与研究目标")

bg_items = [
    "大模型推理开销增长快，单节点难以满足时延与吞吐需求",
    "分布式推理是自然选择，但多节点协同引入信任问题",
    "需要一种「不泄露细节即可验证计算正确」的机制",
    "零知识证明 (ZKP) 恰好满足这一需求",
]
add_bullet_list(s, Inches(0.6), Inches(1.5), Inches(5.5), Inches(4),
                bg_items, title="研究背景")

goals = [
    "构建分布式推理原型系统",
    "引入零知识证明验证推理正确性",
    "设计低开销优化策略降低验证成本",
]
add_bullet_list(s, Inches(7), Inches(1.5), Inches(5.5), Inches(4),
                goals, title="任务书三项目标", title_color=ACCENT)

add_conclusion(s, '课题围绕"分布式推理 + 可验证性 + 低开销"三条主线展开。')


# ═══════════════════════ P3 进展概况 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "当前项目进展概况")

done = [
    "Master + Worker 分布式原型已跑通",
    "2/4/8 切片配置均可正常运行",
    "实现 proof-bound output",
    "完成 6 组实验并生成图表",
]
wip = [
    "统一实验管线与指标口径",
    "收紧系统边界表述",
    "整理答辩与论文材料",
]
nxt = [
    "用主系统重跑关键实验",
    "补充缺失的攻击场景",
    "推进论文写作",
]

col_w = Inches(3.8)
add_bullet_list(s, Inches(0.4), Inches(1.5), col_w, Inches(4.5),
                done, title="已完成 ✅", title_color=GREEN)
add_bullet_list(s, Inches(4.6), Inches(1.5), col_w, Inches(4.5),
                wip, title="正在完善 🔧", title_color=ORANGE)
add_bullet_list(s, Inches(8.8), Inches(1.5), col_w, Inches(4.5),
                nxt, title="下一步 📋", title_color=ACCENT)

add_conclusion(s, "原型已完成，当前重点是实验收口与论文推进。")


# ═══════════════════════ P4 系统架构 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "系统总体架构设计")

arch_items = [
    "Master：统一调度、汇总结果、执行多层验证",
    "Worker：局部子模型推理 + 按角色生成证明",
    "proof 节点 → /infer（完整 ZKP 证明）",
    "light 节点 → /infer_light（推理+哈希）",
    "随机挑战 → /re_prove（概率性抽查）",
]
add_bullet_list(s, Inches(0.6), Inches(1.5), Inches(5), Inches(4.5), arch_items)

# 架构文字图示
arch_text = [
    "Input → Master (调度+校验)",
    "  ├── Worker 1 (/infer) → Slice 1",
    "  ├── Worker 2 (/infer_light) → Slice 2",
    "  └── Worker N (/infer) → Slice N",
    "",
    "校验路径:",
    "  L1: SHA-256 哈希校验",
    "  L2: EZKL proof + proof-bound output",
    "  L3: 跨节点哈希链",
    "  Challenge: /re_prove 随机挑战",
]
add_text_box(s, Inches(6.2), Inches(1.5), Inches(6.5), Inches(4.5),
             arch_text, font_size=14, color=BLUE)

add_conclusion(s, "采用 Master-Worker 架构，将推理流与验证流分层解耦。")


# ═══════════════════════ P5 核心验证机制 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "分层可验证机制设计")

layers = [
    ("L1：外部完整性检查", [
        "SHA-256 哈希校验输入/输出",
        "开销极低，快速发现非恶意故障"
    ]),
    ("L2：零知识证明验证", [
        "Halo2 电路中生成 EZKL 证明",
        "Master 独立验证 proof",
        "proof-bound output：输出从 proof 公开实例提取"
    ]),
    ("L3：跨节点哈希链", [
        "前一节点 hash_out == 后一节点 hash_in",
        "检测链路传输异常"
    ]),
    ("随机挑战", [
        "对 light 节点发起 /re_prove",
        "提供概率性威慑"
    ]),
]

y = Inches(1.5)
for title, items in layers:
    add_bullet_list(s, Inches(0.6), y, Inches(12), Inches(1.2),
                    items, title=title, font_size=14)
    y += Inches(1.2)

add_conclusion(s, "分层验证：proof 节点提供密码学保证，light 节点提供轻量完整性检查。")


# ═══════════════════════ P6 模块完成情况 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "模块划分与当前实现情况")

from pptx.util import Inches as In

# 表格
rows, cols = 6, 3
tbl_shape = s.shapes.add_table(rows, cols,
                                Inches(1.5), Inches(1.6),
                                Inches(10), Inches(3))
tbl = tbl_shape.table

headers = ["模块", "功能", "状态"]
data = [
    ["master.py / worker.py", "调度、推理、验证", "✅"],
    ["utils.py", "EZKL 与哈希工具", "✅"],
    ["run_stage2.py", "主系统验证", "✅"],
    ["run_*experiments.py", "性能与安全实验", "✅"],
    ["test_core_semantics.py", "18 项回归测试", "✅"],
]

for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for ri, row_data in enumerate(data):
    for ci, val in enumerate(row_data):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
            p.alignment = PP_ALIGN.CENTER if ci == 2 else PP_ALIGN.LEFT

notes = [
    "主系统包含完整三层校验、随机挑战与 proof-bound output",
    "部分实验脚本为简化评估管线，用于性能趋势分析",
]
add_bullet_list(s, Inches(1.5), Inches(5), Inches(10), Inches(1.2),
                notes, font_size=14, color=GRAY)

add_conclusion(s, "代码模块体系完整，具备支撑后续实验与论文写作的基础。")


# ═══════════════════════ P7 运行证据 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "主系统运行情况与阶段性证据")

obs = [
    "2/4/8 切片均成功完成端到端推理",
    "证明生成占时延约 85%，是主要瓶颈",
    "篡改被 proof-bound 机制在源头替换",
]
add_bullet_list(s, Inches(0.5), Inches(1.5), Inches(4.5), Inches(3), obs, font_size=15)

img = fig('p07_main_system_latency_breakdown.png')
add_image_safe(s, img, Inches(5.5), Inches(1.3), height=Inches(4.8))

img2 = fig('p07_side_proof_bound_prevention.png')
add_image_safe(s, img2, Inches(9.5), Inches(4.5), height=Inches(2.5))

add_conclusion(s, "系统在多种切片配置下均可稳定运行，proof-bound 机制有效预防响应层篡改。")


# ═══════════════════════ P8 选择性验证 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "选择性验证的低开销效果")

sel_notes = [
    "verify_ratio：用户设定的期望验证比例",
    "actual_proof_fraction：edge-cover 调整后的实际覆盖率",
    "两者不一定相等",
    "",
    "端到端开销降低约 36%~42%",
]
add_bullet_list(s, Inches(0.5), Inches(1.5), Inches(4), Inches(3.5),
                sel_notes, font_size=14)

img = fig('p08_main_selective_verification.png')
add_image_safe(s, img, Inches(4.8), Inches(1.3), height=Inches(4.2))

img2 = fig('p08_side_cost_reduction.png')
add_image_safe(s, img2, Inches(0.5), Inches(4.5), height=Inches(2))

add_conclusion(s, "选择性验证在保持边覆盖约束的前提下，端到端开销降低约 36%~42%。")


# ═══════════════════════ P9 攻击检测 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "攻击检测实验结果")

img = fig('p09_main_attack_handling.png')
add_image_safe(s, img, Inches(3), Inches(1.3), height=Inches(4.2))

atk_notes = [
    "攻击类型：响应层篡改（计算后篡改返回值）",
    "四种变体：tamper / skip / random / replay",
    "恶意节点位于 proof 节点（首尾必选）",
    "",
    "⚠ 结论仅在当前攻击模型下成立",
    "⚠ 不应泛化为「检测所有恶意行为」",
]
add_bullet_list(s, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5),
                atk_notes, font_size=13, color=BLACK)

add_conclusion(s, "在当前攻击模型下，四类响应层篡改均被 proof-bound 机制成功预防。")


# ═══════════════════════ P10 可见性模式 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "不同可见性模式的开销对比")

img = fig('p10_main_visibility_time.png')
add_image_safe(s, img, Inches(0.3), Inches(1.3), height=Inches(4.3))

img2 = fig('p10_side_visibility_size.png')
add_image_safe(s, img2, Inches(6.8), Inches(1.3), height=Inches(4.3))

vis_notes = [
    "hashed 模式约 1.6× 时间开销（Poseidon 哈希电路约束）",
    "private 与 all_public 基本持平",
    "hashed 的 proof 大小增加约 18%",
]
add_bullet_list(s, Inches(0.5), Inches(5.5), Inches(12), Inches(1),
                vis_notes, font_size=14)

add_conclusion(s, "三种可见性模式的开销差异已可定量比较，为隐私-性能权衡提供数据支撑。")


# ═══════════════════════ P11 问题与边界 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "当前存在的问题与系统边界")

problems = [
    "部分实验使用简化管线，尚未与主系统完全统一",
    "verify_ratio 与 actual_proof_fraction 需明确区分",
    "缺少 light 节点被攻击的实验数据",
    "模型规模较小（480 参数 FC），外部效度有限",
]
add_bullet_list(s, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5),
                problems, title="当前待解决问题", title_color=ORANGE)

bounds = [
    "Master 为可信假设，被攻破则全线失效",
    "light 节点非强密码学约束",
    "跨节点中间数据为明文传输",
    "当前为单机多进程模拟",
    "攻击模型仅覆盖响应层篡改",
]
add_bullet_list(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.5),
                bounds, title="系统边界（需诚实面对）", title_color=ORANGE)

add_conclusion(s, "原型已可运行，但需继续收紧实验口径并诚实界定安全边界。")


# ═══════════════════════ P12 下一步计划 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "下一步工作计划")

near = [
    "用 Master 完整逻辑重跑关键实验",
    "补充 light 节点被攻击的场景",
    "统一 actual_proof_fraction 口径",
]
mid = [
    "收紧论文措辞与边界说明",
    "完善性能评估与安全性分析",
    "视情况补充更大模型实验",
]
final = [
    "形成完整实验矩阵",
    "完成毕业论文",
    "完成终期答辩材料",
]

add_bullet_list(s, Inches(0.5), Inches(1.5), Inches(3.8), Inches(4.5),
                near, title="近期（1-2 周）", title_color=GREEN)
add_bullet_list(s, Inches(4.8), Inches(1.5), Inches(3.8), Inches(4.5),
                mid, title="中期后（2-4 周）", title_color=ACCENT)
add_bullet_list(s, Inches(9), Inches(1.5), Inches(3.8), Inches(4.5),
                final, title="最终目标", title_color=BLUE)

add_conclusion(s, "后续核心任务：统一实验口径、补齐关键结果、推进论文收口。")


# ═══════════════════════ P13 总结页 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BLUE)

add_text_box(s, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             ["阶段性总结"],
             font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

summary = [
    "✅ 分布式推理原型系统（Master + Worker，2/4/8 切片）",
    "✅ EZKL (Halo2/PLONK/KZG) 零知识证明接入与验证",
    "✅ proof-bound output（proof 与推理输出数学绑定）",
    "✅ 6 组阶段性实验（性能、选择性验证、攻击、可见性等）",
]
add_text_box(s, Inches(1.5), Inches(2), Inches(10), Inches(2.5),
             summary, font_size=17, color=WHITE)

next_items = [
    "→ 统一实验管线与指标口径",
    "→ 收紧系统边界与论文表述",
    "→ 完成毕业论文与终期答辩材料",
]
add_text_box(s, Inches(1.5), Inches(4.5), Inches(10), Inches(1.5),
             next_items, font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB))

add_text_box(s, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
             ["课题主线已落地为可运行原型，后续重点是实验收口与论文推进。"],
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════ 备份页 A ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "[备份] 切片逻辑一致性验证 (P4)", color=GRAY)

img = fig('bkA_slice_logic_consistency.png')
add_image_safe(s, img, Inches(2.5), Inches(1.5), height=Inches(4.5))

add_text_box(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             ["注：仅验证 PyTorch 切片一致性，非 ONNXRuntime/EZKL 量化路径保真度"],
             font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════ 备份页 B ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "[备份] 三类完整性检查机制对比 (P6)", color=GRAY)

img = fig('bkB_integrity_mechanism_cost.png')
add_image_safe(s, img, Inches(2.5), Inches(1.5), height=Inches(4.5))

add_text_box(s, Inches(0.5), Inches(6.1), Inches(12), Inches(0.5),
             ["仅展示正常模式下三种机制的 proof 时间对比，不应表述为「完整 ZK linking 实证」"],
             font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════ 备份页 C ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_title_bar(s, "[备份] 补充实验图表", color=GRAY)

img = fig('bk_slice8_per_slice_proof_verify.png')
add_image_safe(s, img, Inches(0.5), Inches(1.5), height=Inches(4.5))

img = fig('bk_system_throughput.png')
add_image_safe(s, img, Inches(7), Inches(1.5), height=Inches(4.5))

add_text_box(s, Inches(0.5), Inches(6.2), Inches(5.5), Inches(0.4),
             ["8 切片逐 slice proof/verify 时间"],
             font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(7), Inches(6.2), Inches(5.5), Inches(0.4),
             ["吞吐量随切片数变化"],
             font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════ 感谢页 ═══════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BLUE)

add_text_box(s, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
             ["感谢观看"],
             font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
             ["THANK YOU FOR WATCHING"],
             font_size=20, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(5.5), Inches(11), Inches(1),
             ["答辩人：武垚乐    |    2026 年 3 月"],
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ── 保存 ────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"✅ PPT 已生成: {os.path.abspath(OUT_PATH)}")
print(f"   共 {len(prs.slides)} 页（含 3 备份页 + 感谢页）")
